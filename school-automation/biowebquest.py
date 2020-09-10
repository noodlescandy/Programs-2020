from selenium import webdriver
from selenium.common.exceptions import NoSuchElementException
from selenium.webdriver.firefox.options import Options
from selenium.webdriver.common.keys import Keys
import time

slides = [
["Biology Vocab Webquest", "Timothy Carlson", 0]
]

options = Options()
#options.headless = True
options.headless = False
driver = webdriver.Firefox(options=options)

driver.get("https://biologydictionary.net/")
count = 0
with open('vocablistout.txt', 'r') as vocabFile:
    while(True):
        count += 1
        line = vocabFile.readline()
        if not line:
            break
        searchbox = driver.find_element_by_xpath('/html/body/div[6]/div[2]/div/div[2]/div[1]/div/div[2]/div/div[1]/input')
        searchbox.send_keys(line + Keys.RETURN)
        #searchbutton = driver.find_element_by_xpath('/html/body/div[1]/div/header/div/div[4]/form/div[1]/div[2]/span/svg')
        #searchbutton.click()
        time.sleep(0.5)
        #definition = driver.find_element_by_xpath('/html/body/div[1]/div/div[5]/div[1]/div[1]/div[4]/div[2]/div[1]/span[1]/div/div[1]/span[2]/span').text
        definition = driver.find_element_by_xpath('/html/body/div[1]/div/div[5]/div[1]/div[1]/div[4]/div[2]').text
        print(definition)
        slides.insert(count, [line, definition, 3])
        homepagebutton = driver.find_element_by_xpath('/html/body/div[1]/div/header/div/a[1]/img')









from pptx import Presentation
import os, subprocess, sys

# Create a powerpoint file
prs = Presentation()

class Slide:

	def __init__(self, data):
		# data is a list with title, subtitle and layout num
		print(data)
		self.layout = prs.slide_layouts[data[2]]
		self.slide = prs.slides.add_slide(self.layout)
		self.title = self.slide.shapes.title
		self.title.text = data[0]
		self.subtitle = self.slide.placeholders[1]
		self.subtitle.text = data[1]


# title / subtitle / layout num


def save(name):
    prs.save(name)
    # Open the file on the operating system
    #os.startfile(name) #DOES NOT WORK ON LINUX
    #subprocess.call([opener, name])
    subprocess.call(["xdg-open", name])

for s in slides:
	Slide(s)

save("webquest.pptx")
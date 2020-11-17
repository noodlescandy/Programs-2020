import time, json
from PyDictionary import PyDictionary

slides = [
["Vocab Webquest", "Timothy Carlson", 0]
]

dictionary = PyDictionary()

count = 0
with open('vocablistout.txt', 'r') as vocabFile:
    while(True):
        count += 1
        line = vocabFile.readline()
        if not line:
            break
        definitionDict = dictionary.meaning(line)
        if not(definitionDict is None):
            definition = json.dumps(definitionDict, indent=0)
            definition = definition.replace("{", "")
            definition = definition.replace("\"", "")
            definition = definition.replace(":", "")
            definition = definition.replace(",", "")
            definition = definition.replace("[", "")
            definition = definition.replace("]", "")
            definition = definition.replace("}", "")
            definition = definition.replace("Noun", "")
            definition = definition.replace("Adjective", "")
            definition = definition[2:]
            #print(definition)
            slides.insert(count, [line, definition, 3])
        #print(line + "is null")









from pptx import Presentation
import os, subprocess, sys

# Create a powerpoint file
prs = Presentation()

class Slide:

	def __init__(self, data):
		# data is a list with title, subtitle and layout num
		print(data)
		self.layout = prs.slide_layouts[data[2]]
		self.slide = prs.slides.add_slide(self.layout)
		self.title = self.slide.shapes.title
		self.title.text = data[0]
		self.subtitle = self.slide.placeholders[1]
		self.subtitle.text = data[1]


# title / subtitle / layout num


def save(name):
    prs.save(name)
    # Open the file on the operating system
    #os.startfile(name) #DOES NOT WORK ON LINUX
    #subprocess.call([opener, name])
    subprocess.call(["xdg-open", name])

for s in slides:
	Slide(s)

save("webquest.pptx")

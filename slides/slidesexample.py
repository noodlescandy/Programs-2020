from pptx import Presentation
import os, subprocess, sys

# Create a powerpoint file
prs = Presentation()

class Slide:

	def __init__(self, data):
		# data is a list with title, subtitle and layout num
		print(data)
		self.layout = prs.slide_layouts[data[2]]
		self.slide = prs.slides.add_slide(self.layout)
		self.title = self.slide.shapes.title
		self.title.text = data[0]
		self.subtitle = self.slide.placeholders[1]
		self.subtitle.text = data[1]


# title / subtitle / layout num
slides = [
["Hello people", "We are great", 1], ["Test", "Bazinga", 2]
]

def save(name):
    prs.save(name)
    # Open the file on the operating system
    #os.startfile(name) #DOES NOT WORK ON LINUX
    #subprocess.call([opener, name])
    subprocess.call(["xdg-open", name])

for s in slides:
	Slide(s)

save("prova.pptx")
